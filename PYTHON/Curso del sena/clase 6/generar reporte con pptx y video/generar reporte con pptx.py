from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from reportlab.lib.units import inch
from reportlab.platypus import Table, TableStyle, Paragraph
from reportlab.lib import colors
from reportlab.lib.styles import ParagraphStyle
from pathlib import Path
import pandas as pd
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = Path(__file__).resolve().parent
OUTPUT_DIR = BASE / "resultados"
OUTPUT_DIR.mkdir(exist_ok=True)  # Crear carpeta de resultados si no existe
LOGO_PATH = BASE / "logo_ECCI.jpg"
VIDEO_PATH = BASE / "2025-11-03 20-23-33.mp4"


def graficos(resultados):
    # Gráfico: promedios por programa
    df_programa = resultados['por_programa']
    plt.figure()
    df_programa['promedio'].plot(kind='bar', title="Promedio por Programa")
    plt.ylabel("Promedio")
    plt.tight_layout()
    plt.savefig(OUTPUT_DIR / "grafico_programa.png")
    plt.close()

    # Gráfico: asignaturas mayor variabilidad
    df_asig = resultados['asignaturas_mayor_variabilidad']
    plt.figure()
    df_asig['desviacion'].plot(kind='bar', title="Asignaturas con Mayor Variabilidad")
    plt.ylabel("Desviación Estándar")
    plt.tight_layout()
    plt.savefig(OUTPUT_DIR / "grafico_variabilidad.png")
    plt.close()


def tabla_pdf(pdf, df, x, y, max_width=520):
    # Estilo para salto de línea en texto largo
    style = ParagraphStyle(name="Tabla", fontSize=8, alignment=1)

    # Convertir fila por fila a Paragraph para que se vea el texto completo
    data = [[Paragraph(str(col), style) for col in df.columns]]
    for row in df.round(2).values.tolist():
        data.append([Paragraph(str(cell), style) for cell in row])

    # Ajuste uniforme de ancho
    col_widths = [max_width / len(df.columns)] * len(df.columns)

    table = Table(data, colWidths=col_widths)

    table.setStyle(TableStyle([
        ('BACKGROUND',(0,0),(-1,0), colors.HexColor("#003366")),
        ('TEXTCOLOR',(0,0),(-1,0), colors.white),
        ('GRID',(0,0),(-1,-1), 0.5, colors.black),
        ('FONTNAME',(0,0),(-1,0), 'Helvetica-Bold'),
        ('ALIGN',(0,0),(-1,-1),'CENTER'),
        ('VALIGN',(0,0),(-1,-1),'MIDDLE'),
        ('BOTTOMPADDING',(0,0),(-1,0), 6),
        ('LEFTPADDING',(0,0),(-1,-1), 2),
        ('RIGHTPADDING',(0,0),(-1,-1), 2),
    ]))

    w, h = table.wrapOn(pdf, x, y)
    table.drawOn(pdf, x, y - h)


def generar_pdf(resultados):
    graficos(resultados)

    pdf_path = OUTPUT_DIR / "Informe_Analitica_Ejecutivo.pdf"
    pdf = canvas.Canvas(str(pdf_path), pagesize=A4)
    width, height = A4

    # ✅ PORTADA
    pdf.setFont("Helvetica-Bold", 20)
    pdf.drawString(80, height - 100, "Informe Ejecutivo de Analítica Institucional")
    pdf.setFont("Helvetica", 12)
    pdf.drawString(80, height - 130, "Rendimiento Académico - Universidad ECCI")
    pdf.drawString(80, height - 150, "Generado automáticamente con Python")
    pdf.showPage()

    # ✅ GRÁFICO 1
    pdf.setFont("Helvetica-Bold", 14)
    pdf.drawString(80, height - 50, "Promedio por Programa")
    pdf.drawImage(OUTPUT_DIR / "grafico_programa.png", 50, height - 400, width=500, preserveAspectRatio=True)
    pdf.showPage()

    # ✅ GRÁFICO 2
    pdf.setFont("Helvetica-Bold", 14)
    pdf.drawString(80, height - 50, "Asignaturas con Mayor Variabilidad")
    pdf.drawImage(OUTPUT_DIR / "grafico_variabilidad.png", 50, height - 400, width=500, preserveAspectRatio=True)
    pdf.showPage()

    # ✅ TABLA 1 - Programas (índice convertido a columna)
    pdf.setFont("Helvetica-Bold", 14)
    pdf.drawString(50, height - 50, "Resumen por Programa")
    df_prog = resultados['por_programa'].reset_index().rename(columns={'programa': 'Programa'})
    tabla_pdf(pdf, df_prog, 50, height - 100)
    pdf.showPage()

    # ✅ TABLA 2 - Asignaturas (índice convertido a columna)
    pdf.setFont("Helvetica-Bold", 14)
    pdf.drawString(50, height - 50, "Asignaturas con Mayor Variabilidad")
    df_asig = resultados['asignaturas_mayor_variabilidad'].reset_index().rename(columns={'asignatura': 'Asignatura'})
    tabla_pdf(pdf, df_asig, 50, height - 100)
    pdf.showPage()

    # ✅ CONCLUSIONES
    riesgo = resultados['porcentaje_en_riesgo']
    pdf.setFont("Helvetica-Bold", 16)
    pdf.drawString(50, height - 80, "Conclusiones del Análisis")
    pdf.setFont("Helvetica", 12)
    pdf.drawString(50, height - 120, f"- {riesgo}% de estudiantes están en riesgo académico (promedio < 3.0).")
    pdf.drawString(50, height - 150, "- Las asignaturas con mayor variabilidad requieren atención prioritaria.")
    pdf.drawString(50, height - 180, "- Se recomienda fortalecer tutorías y seguimiento a grupos críticos.")

    pdf.save()
    print(f"✅ PDF generado correctamente: {pdf_path}")


def generar_pptx(resultados):
    """Genera una presentación PPTX con logo y video en la primera diapositiva"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ✅ PRIMERA DIAPOSITIVA: Logo y Video
    slide_layout = prs.slide_layouts[6]  # Diseño en blanco
    slide = prs.slides.add_slide(slide_layout)

    # Agregar logo
    if LOGO_PATH.exists():
        logo_left = Inches(1)
        logo_top = Inches(0.5)
        logo_width = Inches(2)
        logo_height = Inches(1.5)
        slide.shapes.add_picture(str(LOGO_PATH), logo_left, logo_top, logo_width, logo_height)
        print(f"✅ Logo agregado: {LOGO_PATH}")
    else:
        print(f"⚠️ Logo no encontrado en: {LOGO_PATH}")

    # Agregar video
    if VIDEO_PATH.exists():
        video_left = Inches(4)
        video_top = Inches(1.5)
        video_width = Inches(5.5)
        video_height = Inches(4)
        try:
            slide.shapes.add_movie(
                str(VIDEO_PATH),
                video_left, video_top, video_width, video_height,
                poster_frame_image=None,
                mime_type='video/mp4'
            )
            print(f"✅ Video agregado: {VIDEO_PATH}")
        except Exception as e:
            print(f"⚠️ Error al agregar video: {e}")
            # Si falla, agregar un placeholder de texto
            textbox = slide.shapes.add_textbox(Inches(4), Inches(1.5), Inches(5.5), Inches(4))
            text_frame = textbox.text_frame
            text_frame.text = f"Video: {VIDEO_PATH.name}"
    else:
        print(f"⚠️ Video no encontrado en: {VIDEO_PATH}")

    # Título en la primera diapositiva
    title_left = Inches(1)
    title_top = Inches(0.2)
    title_width = Inches(8)
    title_height = Inches(0.5)
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.text = "Informe Ejecutivo de Analítica Institucional"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # ✅ DIAPOSITIVA 2: Gráfico Promedio por Programa
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Promedio por Programa"
    title_frame.paragraphs[0].font.size = Pt(20)
    title_frame.paragraphs[0].font.bold = True
    
    if (OUTPUT_DIR / "grafico_programa.png").exists():
        slide.shapes.add_picture(
            str(OUTPUT_DIR / "grafico_programa.png"),
            Inches(1), Inches(1.5), Inches(8), Inches(5)
        )

    # ✅ DIAPOSITIVA 3: Gráfico Variabilidad
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Asignaturas con Mayor Variabilidad"
    title_frame.paragraphs[0].font.size = Pt(20)
    title_frame.paragraphs[0].font.bold = True
    
    if (OUTPUT_DIR / "grafico_variabilidad.png").exists():
        slide.shapes.add_picture(
            str(OUTPUT_DIR / "grafico_variabilidad.png"),
            Inches(1), Inches(1.5), Inches(8), Inches(5)
        )

    # ✅ DIAPOSITIVA 4: Tabla Resumen por Programa
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Resumen por Programa"
    title_frame.paragraphs[0].font.size = Pt(20)
    title_frame.paragraphs[0].font.bold = True

    df_prog = resultados['por_programa'].reset_index().rename(columns={'programa': 'Programa'})
    
    # Crear tabla
    rows = len(df_prog) + 1
    cols = len(df_prog.columns)
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Encabezados
    for col_idx, col_name in enumerate(df_prog.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Datos
    for row_idx, (_, row) in enumerate(df_prog.iterrows(), start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(round(value, 2))
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ✅ DIAPOSITIVA 5: Tabla Asignaturas con Mayor Variabilidad
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Asignaturas con Mayor Variabilidad"
    title_frame.paragraphs[0].font.size = Pt(20)
    title_frame.paragraphs[0].font.bold = True

    df_asig = resultados['asignaturas_mayor_variabilidad'].reset_index().rename(columns={'asignatura': 'Asignatura'})
    
    rows = len(df_asig) + 1
    cols = len(df_asig.columns)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Encabezados
    for col_idx, col_name in enumerate(df_asig.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Datos
    for row_idx, (_, row) in enumerate(df_asig.iterrows(), start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(round(value, 2))
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ✅ DIAPOSITIVA 6: Conclusiones
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Conclusiones del Análisis"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    riesgo = resultados['porcentaje_en_riesgo']
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    p1 = content_frame.paragraphs[0]
    p1.text = f"• {riesgo}% de estudiantes están en riesgo académico (promedio < 3.0)."
    p1.font.size = Pt(14)
    p1.space_after = Pt(12)
    
    p2 = content_frame.add_paragraph()
    p2.text = "• Las asignaturas con mayor variabilidad requieren atención prioritaria."
    p2.font.size = Pt(14)
    p2.space_after = Pt(12)
    
    p3 = content_frame.add_paragraph()
    p3.text = "• Se recomienda fortalecer tutorías y seguimiento a grupos críticos."
    p3.font.size = Pt(14)

    pptx_path = OUTPUT_DIR / "Informe_Analitica_Ejecutivo.pptx"
    prs.save(str(pptx_path))
    print(f"✅ PPTX generado correctamente: {pptx_path}")


if __name__ == "__main__":
    from analysis_institucional import cargar_fuentes, limpiar_datos, integrar, calcular_estadisticas
    
    df_enroll, df_grades, df_subj = cargar_fuentes()
    df_enroll, df_grades_clean = limpiar_datos(df_enroll, df_grades)
    df_integrado = integrar(df_enroll, df_grades_clean, df_subj)
    resultados = calcular_estadisticas(df_integrado)

    generar_pdf(resultados)
    generar_pptx(resultados)


